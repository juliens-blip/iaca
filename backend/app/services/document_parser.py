"""
Service d'extraction de texte depuis des documents PDF, PPTX et DOCX.
"""

import asyncio
import os


async def parse_document(file_path: str) -> str:
    """
    Extrait le texte brut d'un fichier PDF, PPTX ou DOCX.

    Args:
        file_path: Chemin absolu vers le fichier a parser.

    Returns:
        Le texte brut extrait du document.

    Raises:
        ValueError: Si le type de fichier n'est pas supporte.
        FileNotFoundError: Si le fichier n'existe pas.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Fichier non trouve: {file_path}")

    extension = os.path.splitext(file_path)[1].lower()

    if extension == ".pdf":
        return await _extract_pdf_async(file_path)
    elif extension == ".pptx":
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, _extract_pptx, file_path)
    elif extension == ".docx":
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, _extract_docx, file_path)
    else:
        raise ValueError(f"Type de fichier non supporte: {extension}")


async def _extract_pdf_async(file_path: str) -> str:
    """Extrait le contenu d'un PDF: tente marker (markdown structuré), puis fitz en fallback."""
    from backend.app.services.marker_parser import parse_pdf_with_marker
    return await parse_pdf_with_marker(file_path)


def _extract_pdf(file_path: str) -> str:
    """Extrait le texte d'un fichier PDF via pymupdf (fitz) — conservé pour usage direct."""
    import fitz  # pymupdf

    text_parts: list[str] = []
    with fitz.open(file_path) as doc:
        for page in doc:
            page_text = page.get_text()
            if page_text.strip():
                text_parts.append(page_text)
    return "\n\n".join(text_parts)


def _extract_pptx(file_path: str) -> str:
    """Extrait le texte d'un fichier PowerPoint PPTX."""
    from pptx import Presentation

    text_parts: list[str] = []
    prs = Presentation(file_path)
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        slide_texts.append(para_text)
        if slide_texts:
            text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))
    return "\n\n".join(text_parts)


def _extract_docx(file_path: str) -> str:
    """Extrait le texte d'un fichier Word DOCX."""
    from docx import Document

    text_parts: list[str] = []
    doc = Document(file_path)
    for paragraph in doc.paragraphs:
        para_text = paragraph.text.strip()
        if para_text:
            text_parts.append(para_text)
    return "\n\n".join(text_parts)
