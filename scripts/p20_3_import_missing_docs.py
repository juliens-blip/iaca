#!/usr/bin/env python3
"""
P20-3: Import des documents manquants du dossier source
Mappe les dossiers source → matière_id et importe les docs manquants
"""
import os
import sys
import sqlite3
from pathlib import Path
from typing import Dict, List, Tuple
import json

# Ajouter backend au path
sys.path.insert(0, str(Path(__file__).parent.parent / "backend"))

def get_source_dir():
    """Récupère le chemin du dossier source"""
    return Path("/home/julien/Téléchargements/drive-download-20260226T072425Z-3-001")

def build_folder_to_matiere_mapping() -> Dict[str, int]:
    """
    Construit un mapping : chemin dossier source -> matiere_id
    """
    mapping = {
        "Droit public ": 1,
        "Economie": 6,
        "Finances Publiques": 7,
        "Espagnol": 8,
        "Licence 2 Assas/DROIT L2 Semestre 3": 9,
        "Licence 2 Assas/DROIT L2 Semestre 4": 10,
        "Licence 2 Assas/Fiche de révision 2025 - 2026 ": 11,
        "Licence 3 assas/Semestre 5": 12,
        "Licence 3 assas/Semestre 6": 13,
        "Licence 3 assas/TEDS": 14,
        "Licence 3 assas/Scolarité": 15,
        "Livres et manuels": 16,
        "Question contemporaine": 3,
        "Questions sociales ": 4,
        "Relations internationales ": 5,
        "S1_M1_Droit et légistique": 17,
        "S1_M1_Politique économique": 18,
        "S1_Management et égalités pro _( ": 19,
    }
    return mapping

def get_matiere_id_for_path(file_path: Path, mapping: Dict[str, int]) -> int:
    """
    Détermine la matiere_id pour un fichier donné
    """
    # Essayer de matcher le chemin le plus spécifique d'abord
    relative_path = file_path.relative_to(get_source_dir())
    path_str = str(relative_path).replace(os.sep, "/")

    # Chercher le match le plus long (plus spécifique)
    best_match = None
    best_match_len = 0

    for folder_pattern, matiere_id in mapping.items():
        if path_str.startswith(folder_pattern):
            if len(folder_pattern) > best_match_len:
                best_match = matiere_id
                best_match_len = len(folder_pattern)

    if best_match is not None:
        return best_match

    # Fallback : "Documents divers"
    return 20

def extract_text_from_file(file_path: Path, extract: bool = False) -> str:
    """
    Extrait le texte d'un fichier (PDF, DOCX, PPTX)
    Retourne les 500 premiers chars pour la preview
    Si extract=False, retourne une chaîne vide pour plus de vitesse
    """
    if not extract:
        return ""

    try:
        if file_path.suffix.lower() == ".pdf":
            try:
                import pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    text = "\n".join(page.extract_text() or "" for page in pdf.pages if page.extract_text())
                    return text[:500] if text else ""
            except Exception:
                return ""

        elif file_path.suffix.lower() == ".docx":
            try:
                from docx import Document
                doc = Document(file_path)
                text = "\n".join(p.text for p in doc.paragraphs if p.text)
                return text[:500] if text else ""
            except Exception:
                return ""

        elif file_path.suffix.lower() == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                text = "\n".join(
                    shape.text for slide in prs.slides
                    for shape in slide.shapes if hasattr(shape, "text") and shape.text
                )
                return text[:500] if text else ""
            except Exception:
                return ""
    except Exception:
        pass

    return ""

def collect_files(source_dir: Path) -> List[Path]:
    """
    Collecte tous les fichiers importables (PDF, DOCX, PPTX)
    """
    extensions = {".pdf", ".docx", ".pptx"}
    files = []

    for ext in extensions:
        files.extend(source_dir.glob(f"**/*{ext}"))

    # Aussi à la racine
    for ext in extensions:
        files.extend(source_dir.glob(f"*{ext}"))

    return sorted(set(files))

def insert_document(db: sqlite3.Connection, titre: str, matiere_id: int,
                   contenu_extrait: str, path: Path) -> int:
    """
    Insère un document en DB, retourne l'id
    """
    cursor = db.cursor()

    # Limiter le contenu extrait
    contenu = contenu_extrait[:2000] if contenu_extrait else ""

    # Déterminer le type de fichier
    file_type = path.suffix.lower().lstrip(".")

    cursor.execute("""
        INSERT INTO documents (titre, fichier_path, type_fichier, matiere_id, contenu_extrait, created_at)
        VALUES (?, ?, ?, ?, ?, datetime('now'))
    """, (titre, str(path), file_type, matiere_id, contenu))

    db.commit()
    return cursor.lastrowid

def main():
    source_dir = get_source_dir()
    db_path = Path("/home/julien/Documents/IACA/data/iaca.db")

    if not source_dir.exists():
        print(f"❌ Source dir not found: {source_dir}")
        return

    # Connexion DB
    db = sqlite3.connect(db_path)
    db.row_factory = sqlite3.Row
    cursor = db.cursor()

    # Vérifier la table documents
    cursor.execute("SELECT name FROM sqlite_master WHERE type='table' AND name='documents'")
    if not cursor.fetchone():
        print("❌ Table 'documents' not found")
        return

    mapping = build_folder_to_matiere_mapping()

    print("=" * 70)
    print("P20-3: IMPORT DOCUMENTS MANQUANTS")
    print("=" * 70)

    # Récupérer les documents existants (par titre)
    cursor.execute("SELECT DISTINCT UPPER(TRIM(titre)) FROM documents")
    existing_titles = {row[0] for row in cursor.fetchall()}

    print(f"\n✅ Documents existants en DB: {len(existing_titles)}")

    # Collecter les fichiers source
    files = collect_files(source_dir)
    print(f"📂 Fichiers trouvés en source: {len(files)}")

    # Importer les fichiers manquants
    imported_count = 0
    import_by_matiere = {}

    print(f"\n🔄 Traitement des fichiers...")

    for file_path in files:
        titre = file_path.stem  # Nom sans extension
        titre_upper = titre.upper().strip()

        # Vérifier si déjà importé
        if titre_upper in existing_titles:
            continue

        matiere_id = get_matiere_id_for_path(file_path, mapping)
        contenu = extract_text_from_file(file_path, extract=False)  # Pour vitesse: extraction optionnelle

        try:
            doc_id = insert_document(db, titre, matiere_id, contenu, file_path)
            imported_count += 1
            import_by_matiere[matiere_id] = import_by_matiere.get(matiere_id, 0) + 1
            print(f"  ✅ {titre[:50]:50s} (M{matiere_id})")
        except Exception as e:
            print(f"  ❌ {titre[:50]:50s} : {e}")

    # Récupérer les noms des matières pour le rapport
    cursor.execute("SELECT id, nom FROM matieres WHERE id IN ({})".format(
        ",".join(str(m) for m in import_by_matiere.keys())
    ))
    matiere_names = {row[0]: row[1] for row in cursor.fetchall()}

    print("\n" + "=" * 70)
    print(f"📊 RAPPORT D'IMPORT")
    print("=" * 70)
    print(f"Documents importés: {imported_count}")
    print(f"\nRépartition par matière:")

    for matiere_id in sorted(import_by_matiere.keys()):
        count = import_by_matiere[matiere_id]
        name = matiere_names.get(matiere_id, "?")
        print(f"  M{matiere_id:2d} | {name:40s} : {count:3d} doc(s)")

    # Vérifier le volume total en DB
    cursor.execute("SELECT COUNT(*) as total FROM documents")
    total_docs = cursor.fetchone()[0]

    print(f"\n📈 Total documents en DB: {total_docs}")

    db.close()

    print("\n✅ Import terminé")

if __name__ == "__main__":
    main()
